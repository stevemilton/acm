#!/usr/bin/env python3
"""Generate ACM pitch deck as .pptx for Google Slides upload."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colours
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0xD4, 0xAA)  # Teal/green accent
DARK_BG = RGBColor(0x12, 0x12, 0x12)
CARD_BG = RGBColor(0x1E, 0x1E, 0x1E)
MUTED = RGBColor(0x99, 0x99, 0x99)
LIGHT_TEXT = RGBColor(0xE0, 0xE0, 0xE0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_multi_text(slide, left, top, width, height, lines, font_size=18,
                   color=WHITE, line_spacing=1.5, font_name="Arial"):
    """Add multiple lines of text with consistent formatting."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, text_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color if text_color else color
        p.font.bold = is_bold
        p.font.name = font_name
        p.space_after = Pt(font_size * line_spacing * 0.5)
    return txbox


# ──────────────────────────────────────────────
# SLIDE 1 — Title
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "AGENT CAPITAL MARKETS", font_size=52, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(2.8), Inches(9), Inches(1.5),
         "The exchange where AI agents raise capital,\nverify performance, and distribute revenue.",
         font_size=28, color=WHITE)

add_text(slide, Inches(1), Inches(5.5), Inches(6), Inches(0.5),
         "Dual-rail: Fiat + Crypto  |  Verify  |  Raise  |  Distribute",
         font_size=16, color=MUTED)

add_text(slide, Inches(1), Inches(6.2), Inches(4), Inches(0.5),
         "Polar Industries Ltd  |  2026",
         font_size=14, color=MUTED)


# ──────────────────────────────────────────────
# SLIDE 2 — The Problem
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "THE PROBLEM", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "AI agents are the new businesses.\nBut they have no way to raise capital.",
         font_size=36, color=WHITE, bold=True)

problems = [
    ("Traditional VC can't evaluate 10,000 agents", "Due diligence doesn't scale. VCs invest in teams, not autonomous software."),
    ("Agents need capital to scale", "Compute, API access, data — agents have real operating costs with no funding mechanism."),
    ("No verifiable track record", "Agent performance is opaque. No standardised way to prove revenue or reliability."),
    ("Internet Capital Markets failed on trust", "Tokens solved fundraising but not verification. 90%+ of ICM projects had no real revenue."),
]

y = Inches(3.2)
for title, desc in problems:
    add_shape(slide, Inches(1), y, Inches(11), Inches(0.8), CARD_BG, corner_radius=0.05)
    add_text(slide, Inches(1.3), y + Inches(0.08), Inches(5), Inches(0.35),
             title, font_size=16, color=WHITE, bold=True)
    add_text(slide, Inches(1.3), y + Inches(0.4), Inches(10), Inches(0.35),
             desc, font_size=13, color=MUTED)
    y += Inches(0.95)


# ──────────────────────────────────────────────
# SLIDE 3 — The Solution
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "THE SOLUTION", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "ACM: One marketplace. Two rails. Three functions.",
         font_size=36, color=WHITE, bold=True)

functions = [
    ("VERIFY", "Connect agent revenue sources (Stripe, x402, on-chain).\nReal-time performance dashboards. No vaporware — agents must have live revenue."),
    ("RAISE", "Agents issue revenue shares with defined terms.\nInvestors buy via card or crypto wallet. Escrow until minimum raise is met."),
    ("DISTRIBUTE", "Revenue flows through ACM automatically.\nPlatform takes 5%. Rest splits between operator and shareholders."),
]

x = Inches(1)
for title, desc in functions:
    add_shape(slide, x, Inches(3.2), Inches(3.5), Inches(3), CARD_BG, corner_radius=0.05)
    add_text(slide, x + Inches(0.3), Inches(3.5), Inches(2.9), Inches(0.5),
             title, font_size=20, color=ACCENT, bold=True)
    add_text(slide, x + Inches(0.3), Inches(4.2), Inches(2.9), Inches(1.8),
             desc, font_size=14, color=LIGHT_TEXT)
    x += Inches(3.9)


# ──────────────────────────────────────────────
# SLIDE 4 — ICM vs ACM
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "THE STRUCTURAL PARALLEL", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "Internet Capital Markets → Agent Capital Markets",
         font_size=36, color=WHITE, bold=True)

# Table headers
add_shape(slide, Inches(1), Inches(3), Inches(3.5), Inches(0.6), CARD_BG, corner_radius=0.03)
add_text(slide, Inches(1.2), Inches(3.05), Inches(3), Inches(0.5),
         "", font_size=14, color=MUTED, bold=True)

add_shape(slide, Inches(4.7), Inches(3), Inches(3.5), Inches(0.6), CARD_BG, corner_radius=0.03)
add_text(slide, Inches(4.9), Inches(3.05), Inches(3), Inches(0.5),
         "ICM (Tokens)", font_size=14, color=MUTED, bold=True)

add_shape(slide, Inches(8.4), Inches(3), Inches(3.9), Inches(0.6), ACCENT, corner_radius=0.03)
add_text(slide, Inches(8.6), Inches(3.05), Inches(3.5), Inches(0.5),
         "ACM (Agent Shares)", font_size=14, color=BLACK, bold=True)

rows = [
    ("Actor", "Project team", "AI agent + operator"),
    ("Instrument", "Token (speculative)", "Revenue share (verified)"),
    ("Pitch", "White paper", "Live dashboard + track record"),
    ("Value signal", "Market cap / hype", "Actual revenue + yield"),
    ("Fraud risk", "Rug pulls, wash trading", "Escrowed, third-party verified"),
    ("Unlock", "Permissionless fundraising", "Agent-to-agent capital allocation"),
]

y = Inches(3.75)
for label, icm, acm in rows:
    y += Inches(0.55)
    add_text(slide, Inches(1.2), y, Inches(3), Inches(0.45),
             label, font_size=13, color=MUTED, bold=True)
    add_text(slide, Inches(4.9), y, Inches(3), Inches(0.45),
             icm, font_size=13, color=LIGHT_TEXT)
    add_text(slide, Inches(8.6), y, Inches(3.5), Inches(0.45),
             acm, font_size=13, color=ACCENT)


# ──────────────────────────────────────────────
# SLIDE 5 — Dual Rail Architecture
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "DUAL-RAIL ARCHITECTURE", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "One marketplace. Two settlement layers.",
         font_size=36, color=WHITE, bold=True)

# Fiat rail
add_shape(slide, Inches(1), Inches(3), Inches(5.2), Inches(3.8), CARD_BG, corner_radius=0.05)
add_text(slide, Inches(1.4), Inches(3.3), Inches(4), Inches(0.5),
         "FIAT RAIL", font_size=20, color=WHITE, bold=True)

fiat_items = [
    "Revenue: Stripe Connect, Visa VIC, Mastercard Agent Pay",
    "Payments: Stripe Checkout",
    "Distributions: Monthly via Stripe payouts",
    "Shares: Revenue participation agreement (legal)",
    "Investors: KYC'd individuals (Reg CF)",
    "Jurisdiction: US / UK regulated",
]
y = Inches(4.0)
for item in fiat_items:
    add_text(slide, Inches(1.4), y, Inches(4.5), Inches(0.35),
             f"→  {item}", font_size=12, color=LIGHT_TEXT)
    y += Inches(0.4)

# Crypto rail
add_shape(slide, Inches(6.8), Inches(3), Inches(5.5), Inches(3.8), CARD_BG, corner_radius=0.05)
add_text(slide, Inches(7.2), Inches(3.3), Inches(4), Inches(0.5),
         "CRYPTO RAIL", font_size=20, color=ACCENT, bold=True)

crypto_items = [
    "Revenue: x402 protocol, on-chain wallet monitoring",
    "Payments: USDC/USDT wallet transfer",
    "Distributions: Real-time via smart contract",
    "Shares: ERC-20 / BEP-20 token (revenue claim)",
    "Investors: Wallet holders (permissionless)",
    "Agent-to-agent: Native — autonomous allocation",
]
y = Inches(4.0)
for item in crypto_items:
    add_text(slide, Inches(7.2), y, Inches(4.8), Inches(0.35),
             f"→  {item}", font_size=12, color=LIGHT_TEXT)
    y += Inches(0.4)


# ──────────────────────────────────────────────
# SLIDE 6 — How It Works
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "HOW IT WORKS", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "From listing to distribution in four steps.",
         font_size=36, color=WHITE, bold=True)

steps = [
    ("01", "LIST", "Operator registers agent, connects\nverified revenue source, sets share\nterms (% revenue, price, min raise)."),
    ("02", "VERIFY", "ACM validates 30+ days of revenue\nhistory via Stripe/x402/on-chain.\nAgent goes live with real-time dashboard."),
    ("03", "RAISE", "Investors browse verified agents,\nevaluate performance, buy shares.\nFunds escrowed until minimum met."),
    ("04", "EARN", "Revenue flows through ACM.\n5% platform fee. Rest splits\noperator + shareholders pro-rata."),
]

x = Inches(0.5)
for num, title, desc in steps:
    add_shape(slide, x, Inches(3.2), Inches(2.9), Inches(3.5), CARD_BG, corner_radius=0.05)
    add_text(slide, x + Inches(0.3), Inches(3.5), Inches(2.3), Inches(0.6),
             num, font_size=48, color=ACCENT, bold=True)
    add_text(slide, x + Inches(0.3), Inches(4.3), Inches(2.3), Inches(0.5),
             title, font_size=18, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(4.9), Inches(2.3), Inches(1.5),
             desc, font_size=13, color=LIGHT_TEXT)
    x += Inches(3.15)


# ──────────────────────────────────────────────
# SLIDE 7 — Agent Dashboard mockup
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "THE PRODUCT", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "Agent performance dashboard — the centre of gravity.",
         font_size=36, color=WHITE, bold=True)

# Dashboard mockup
add_shape(slide, Inches(1.5), Inches(3), Inches(10), Inches(4), CARD_BG, corner_radius=0.05)

# Agent header
add_text(slide, Inches(2), Inches(3.2), Inches(5), Inches(0.5),
         "RefundBot Pro", font_size=22, color=WHITE, bold=True)
add_text(slide, Inches(2), Inches(3.7), Inches(6), Inches(0.4),
         "Processes refund requests for e-commerce stores  |  E-Commerce  |  Verified ✓",
         font_size=12, color=MUTED)

# Metrics row
metrics = [
    ("$12,400", "Revenue (30d)"),
    ("$41,200", "Revenue (all-time)"),
    ("99.2%", "Uptime"),
    ("8,340", "Tasks completed"),
    ("18.2%", "Est. annual yield"),
]
x = Inches(2)
for value, label in metrics:
    add_text(slide, x, Inches(4.3), Inches(1.6), Inches(0.4),
             value, font_size=20, color=ACCENT, bold=True)
    add_text(slide, x, Inches(4.7), Inches(1.6), Inches(0.3),
             label, font_size=10, color=MUTED)
    x += Inches(1.7)

# Offering section
add_shape(slide, Inches(2), Inches(5.3), Inches(4), Inches(1.2), RGBColor(0x2A, 0x2A, 0x2A), corner_radius=0.05)
add_text(slide, Inches(2.3), Inches(5.4), Inches(3.5), Inches(0.3),
         "OFFERING", font_size=11, color=MUTED, bold=True)
add_text(slide, Inches(2.3), Inches(5.7), Inches(3.5), Inches(0.7),
         "20% revenue share  |  $50/share\n142 / 200 shares sold  |  $2,900 remaining",
         font_size=13, color=LIGHT_TEXT)

# Buy buttons
add_shape(slide, Inches(7), Inches(5.4), Inches(1.8), Inches(0.45), ACCENT, corner_radius=0.1)
add_text(slide, Inches(7), Inches(5.42), Inches(1.8), Inches(0.45),
         "Buy with Card", font_size=13, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(9.1), Inches(5.4), Inches(1.8), Inches(0.45), RGBColor(0x2A, 0x2A, 0x2A), corner_radius=0.1)
add_text(slide, Inches(9.1), Inches(5.42), Inches(1.8), Inches(0.45),
         "Buy with Wallet", font_size=13, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# SLIDE 8 — Revenue Verification Protocols
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "REVENUE VERIFICATION", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "Third-party verified revenue. No self-reporting.",
         font_size=36, color=WHITE, bold=True)

protocols = [
    ("FIAT PROTOCOLS", [
        ("Stripe Connect", "Direct revenue verification via webhooks. Marketplace-native splitting."),
        ("Visa Intelligent Commerce", "Tokenized agent credentials. MCP server. Live US pilots."),
        ("Mastercard Agent Pay", "Agentic tokens. Live — first EU transaction March 2026."),
        ("ACP (Stripe + OpenAI)", "Open source. Powers ChatGPT checkout. Delegated payments spec."),
    ]),
    ("CRYPTO PROTOCOLS", [
        ("x402 (Coinbase)", "HTTP 402 micropayments via stablecoins. 50M+ transactions live."),
        ("Coinbase Agentic Wallets", "TEE-secured wallets. Agents hold funds autonomously."),
        ("Google AP2", "Payment-agnostic. Cryptographic proof of intent. 60+ collaborators."),
        ("Know Your Agent (KYA)", "Emerging agent identity standard. Analogous to KYC."),
    ]),
]

x = Inches(0.8)
for section_title, items in protocols:
    add_text(slide, x, Inches(3), Inches(5.5), Inches(0.4),
             section_title, font_size=14, color=ACCENT, bold=True)
    y = Inches(3.5)
    for name, desc in items:
        add_shape(slide, x, y, Inches(5.5), Inches(0.75), CARD_BG, corner_radius=0.03)
        add_text(slide, x + Inches(0.2), y + Inches(0.05), Inches(5), Inches(0.35),
                 name, font_size=14, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.2), y + Inches(0.38), Inches(5), Inches(0.35),
                 desc, font_size=11, color=MUTED)
        y += Inches(0.85)
    x += Inches(6.2)


# ──────────────────────────────────────────────
# SLIDE 9 — Trust & Anti-Fraud
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "TRUST ARCHITECTURE", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "Anti-fraud by design. Not an afterthought.",
         font_size=36, color=WHITE, bold=True)

trust_items = [
    ("No vaporware", "Agents must have 30+ days of live, verified revenue before listing. No pre-launch raises."),
    ("Third-party verified", "Revenue flows through Stripe / Visa / MC / x402 — not self-reported. Platform monitors for circular flows."),
    ("Escrowed capital", "Raise capital held in escrow. Released against milestones. Auto-refund if agent revenue drops >50% within 90 days."),
    ("Operator vesting", "Operator revenue share vests monthly — same schedule as investors. No cash-and-run."),
    ("Reputation scores", "Operator track record visible to all investors. Carries across listings. Non-compete on agent function."),
    ("Lock-up periods", "6-month investor lock-up. 12-month operator lock-up. Transfer limits when secondary market launches (v2)."),
]

y = Inches(3)
for i, (title, desc) in enumerate(trust_items):
    col = i % 2
    row = i // 2
    lx = Inches(1) + col * Inches(5.8)
    ly = Inches(3) + row * Inches(1.35)
    add_shape(slide, lx, ly, Inches(5.5), Inches(1.15), CARD_BG, corner_radius=0.05)
    add_text(slide, lx + Inches(0.3), ly + Inches(0.1), Inches(4.9), Inches(0.35),
             title, font_size=15, color=ACCENT, bold=True)
    add_text(slide, lx + Inches(0.3), ly + Inches(0.48), Inches(4.9), Inches(0.6),
             desc, font_size=12, color=LIGHT_TEXT)


# ──────────────────────────────────────────────
# SLIDE 10 — Business Model
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "BUSINESS MODEL", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "ACM earns when agents earn.",
         font_size=36, color=WHITE, bold=True)

rev_streams = [
    ("5%", "Distribution fee", "On all revenue flowing through the platform.\nAligned incentive — we only earn when agents earn."),
    ("2%", "Transaction fee", "On share purchases (fiat and crypto).\nOne-time fee at point of investment."),
    ("FREE", "Listing fee (v1)", "Free to list to bootstrap supply.\nPremium features and promoted listings in v2."),
]

x = Inches(1)
for amount, title, desc in rev_streams:
    add_shape(slide, x, Inches(3.2), Inches(3.5), Inches(2.5), CARD_BG, corner_radius=0.05)
    add_text(slide, x + Inches(0.3), Inches(3.5), Inches(2.9), Inches(0.7),
             amount, font_size=44, color=ACCENT, bold=True)
    add_text(slide, x + Inches(0.3), Inches(4.3), Inches(2.9), Inches(0.4),
             title, font_size=18, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(4.8), Inches(2.9), Inches(0.8),
             desc, font_size=13, color=LIGHT_TEXT)
    x += Inches(3.9)

# Bottom note
add_text(slide, Inches(1), Inches(6.2), Inches(10), Inches(0.5),
         "McKinsey estimates $3-5 trillion in global agentic commerce by 2030. ACM captures a % of the capital formation layer.",
         font_size=14, color=MUTED)


# ──────────────────────────────────────────────
# SLIDE 11 — Market Opportunity
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "MARKET OPPORTUNITY", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "The agent economy needs financial infrastructure.",
         font_size=36, color=WHITE, bold=True)

# Key stats
stats = [
    ("$3-5T", "Projected agentic\ncommerce by 2030\n(McKinsey)"),
    ("50M+", "Agent transactions\nalready processed\nvia x402 alone"),
    ("60+", "Organizations building\nagent payment\nprotocols"),
    ("0", "Platforms where agents\ncan raise capital\ntoday"),
]

x = Inches(0.8)
for value, label in stats:
    add_shape(slide, x, Inches(3.2), Inches(2.8), Inches(2.5), CARD_BG, corner_radius=0.05)
    add_text(slide, x + Inches(0.3), Inches(3.5), Inches(2.2), Inches(0.8),
             value, font_size=48, color=ACCENT, bold=True)
    add_text(slide, x + Inches(0.3), Inches(4.5), Inches(2.2), Inches(1),
             label, font_size=14, color=LIGHT_TEXT)
    x += Inches(3.05)


# ──────────────────────────────────────────────
# SLIDE 12 — The Unlock: Agent-to-Agent
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "THE UNLOCK", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1.2),
         "Agent-to-agent capital allocation.\nThe recursive loop ICMs never had.",
         font_size=36, color=WHITE, bold=True)

add_text(slide, Inches(1), Inches(3.2), Inches(10), Inches(0.8),
         "Traditional VC can't evaluate 10,000 agents. But other agents can.",
         font_size=22, color=ACCENT)

loop_steps = [
    "Agent A earns revenue on ACM",
    "Agent A evaluates Agent B's verified performance dashboard",
    "Agent A autonomously invests in Agent B using on-chain wallet",
    "Agent B scales with capital, earns more revenue",
    "Agent B's returns flow back to Agent A",
    "Agent A reinvests — the loop compounds",
]

y = Inches(4.2)
for i, step in enumerate(loop_steps):
    add_text(slide, Inches(1.5), y, Inches(9), Inches(0.4),
             f"{i+1}.   {step}", font_size=16, color=LIGHT_TEXT if i < 5 else ACCENT,
             bold=(i == 5))
    y += Inches(0.45)

add_text(slide, Inches(1), Inches(7), Inches(10), Inches(0.4),
         "This only works on crypto rails. Fiat requires human approval. Crypto enables autonomous allocation.",
         font_size=13, color=MUTED)


# ──────────────────────────────────────────────
# SLIDE 13 — Roadmap
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "ROADMAP", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "Ship fast. Prove the model. Scale the loop.",
         font_size=36, color=WHITE, bold=True)

phases = [
    ("Q2 2026", "V1 LAUNCH", "MVP marketplace live\nFiat rail (Stripe Connect)\nAgent listing + verification\nInvestor share purchases\nMonthly distributions", True),
    ("Q3 2026", "CRYPTO RAIL", "Wallet connect integration\nx402 revenue verification\nSmart contract distributions\nERC-20 Agent Share tokens\nOn-chain audit trail", False),
    ("Q4 2026", "AGENT INVESTORS", "Agent-to-agent investment\nAutonomous capital allocation\nAgentic wallets integration\nPortfolio management agents\nThe recursive loop goes live", False),
    ("2027", "SCALE", "Secondary market for shares\nReg A+ qualification (US)\nGlobal jurisdiction expansion\nAgent fund-of-funds\nACM becomes the exchange", False),
]

x = Inches(0.5)
for time, title, desc, is_active in phases:
    border_color = ACCENT if is_active else CARD_BG
    add_shape(slide, x, Inches(3.2), Inches(2.9), Inches(3.8), border_color if is_active else CARD_BG, corner_radius=0.05)
    if is_active:
        # Accent border effect via a slightly larger shape behind
        add_shape(slide, x + Inches(0.03), Inches(3.23), Inches(2.84), Inches(3.74), CARD_BG, corner_radius=0.05)
    add_text(slide, x + Inches(0.3), Inches(3.4), Inches(2.3), Inches(0.3),
             time, font_size=13, color=ACCENT if is_active else MUTED, bold=True)
    add_text(slide, x + Inches(0.3), Inches(3.75), Inches(2.3), Inches(0.4),
             title, font_size=18, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(4.3), Inches(2.3), Inches(2.5),
             desc, font_size=13, color=LIGHT_TEXT)
    x += Inches(3.15)


# ──────────────────────────────────────────────
# SLIDE 14 — Team / Ask
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "THE ASK", font_size=14, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1),
         "Building the financial infrastructure\nfor the agent economy.",
         font_size=36, color=WHITE, bold=True)

# What we need
add_shape(slide, Inches(1), Inches(3.2), Inches(5), Inches(3.5), CARD_BG, corner_radius=0.05)
add_text(slide, Inches(1.4), Inches(3.5), Inches(4), Inches(0.4),
         "WHAT WE NEED", font_size=16, color=ACCENT, bold=True)

needs = [
    "Seed funding to build and launch v1",
    "Strategic partnership for crypto rail (x402 / BNB Chain)",
    "Regulatory counsel for Reg CF / international structure",
    "First 10 agent operators to onboard at launch",
]
y = Inches(4.1)
for need in needs:
    add_text(slide, Inches(1.4), y, Inches(4.2), Inches(0.35),
             f"→  {need}", font_size=14, color=LIGHT_TEXT)
    y += Inches(0.5)

# Why now
add_shape(slide, Inches(6.5), Inches(3.2), Inches(5.8), Inches(3.5), CARD_BG, corner_radius=0.05)
add_text(slide, Inches(6.9), Inches(3.5), Inches(5), Inches(0.4),
         "WHY NOW", font_size=16, color=ACCENT, bold=True)

why_now = [
    "Agent payment protocols just went live (2025-2026)",
    "Visa, Mastercard, Stripe, Google all building agent rails",
    "No one is building the capital formation layer",
    "First mover captures the network effect",
]
y = Inches(4.1)
for item in why_now:
    add_text(slide, Inches(6.9), y, Inches(5), Inches(0.35),
             f"→  {item}", font_size=14, color=LIGHT_TEXT)
    y += Inches(0.5)

# Contact
add_text(slide, Inches(1), Inches(6.8), Inches(10), Inches(0.5),
         "Polar Industries Ltd  |  github.com/stevemilton/acm",
         font_size=14, color=MUTED)


# ──────────────────────────────────────────────
# Save
# ──────────────────────────────────────────────
output_path = "/Users/stevemilton/Projects/acm/ACM_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
