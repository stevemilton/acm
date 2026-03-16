#!/usr/bin/env python3
"""Generate ACM pitch deck for Binance Labs — crypto/stablecoin focused."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colours — Binance-sympathetic palette
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BNB_GOLD = RGBColor(0xF0, 0xB9, 0x0B)  # Binance yellow/gold
DARK_BG = RGBColor(0x12, 0x12, 0x12)
CARD_BG = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x88, 0x88, 0x88)
LIGHT_TEXT = RGBColor(0xD0, 0xD0, 0xD0)
GOLD_DIM = RGBColor(0xA0, 0x7A, 0x08)

# Fonts
H1 = "Montserrat"
H2 = "Play"
BODY = "Inter"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=BODY):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def headline(slide, left, top, width, text, size=36):
    return add_text(slide, left, top, width, Inches(1.2), text,
                    font_size=size, color=WHITE, bold=True, font_name=H1)


def subhead(slide, left, top, width, text, size=14, color=BNB_GOLD):
    return add_text(slide, left, top, width, Inches(0.5), text,
                    font_size=size, color=color, bold=True, font_name=H2)


def body(slide, left, top, width, height, text, size=14, color=LIGHT_TEXT):
    return add_text(slide, left, top, width, height, text,
                    font_size=size, color=color, font_name=BODY)


# ──────────────────────────────────────────────
# SLIDE 1 — Title
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(1.2), Inches(10), "FOR BINANCE LABS", size=16)
headline(s, Inches(1), Inches(2), Inches(11),
         "AGENT CAPITAL MARKETS", size=56)
add_text(s, Inches(1), Inches(3.5), Inches(9), Inches(1),
         "The BNB Chain-native exchange where AI agents raise\ncapital through FDUSD revenue shares and distribute\nearnings via smart contracts.",
         font_size=24, color=LIGHT_TEXT, font_name=BODY)

add_text(s, Inches(1), Inches(5.5), Inches(8), Inches(0.4),
         "BNB Chain native  ·  FDUSD settlement  ·  Agent-to-agent capital allocation",
         font_size=15, color=BNB_GOLD, font_name=H2)

add_text(s, Inches(1), Inches(6.5), Inches(4), Inches(0.4),
         "Polar Industries Ltd  |  2026", font_size=13, color=MUTED, font_name=BODY)


# ──────────────────────────────────────────────
# SLIDE 2 — Why Crypto Rails
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "THE THESIS")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "Agents need crypto-native capital rails.\nFiat can't do what's required.")

items = [
    ("Real-time settlement", "Agents operate 24/7. Monthly Stripe payouts are too slow. FDUSD distributions happen in real-time via smart contracts on BNB Chain."),
    ("Autonomous allocation", "The killer feature: agents investing in other agents. This requires programmable money — wallets that agents control without human approval."),
    ("Permissionless & global", "AI agents have no nationality. Fiat rails are jurisdiction-locked. On-chain rails are borderless from day one."),
    ("Verifiable by default", "Every revenue event, every distribution, every investment — on-chain, auditable, immutable. Trust is built into the ledger."),
]

y = Inches(3.3)
for title, desc in items:
    add_shape(s, Inches(1), y, Inches(11), Inches(0.9), CARD_BG, corner_radius=0.04)
    add_text(s, Inches(1.4), y + Inches(0.08), Inches(4), Inches(0.35),
             title, font_size=15, color=BNB_GOLD, bold=True, font_name=H2)
    body(s, Inches(1.4), y + Inches(0.45), Inches(10), Inches(0.4),
         desc, size=12, color=LIGHT_TEXT)
    y += Inches(1.0)


# ──────────────────────────────────────────────
# SLIDE 3 — What ACM Does
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "THE PRODUCT")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "Three functions. One exchange.")

functions = [
    ("VERIFY", "Connect on-chain revenue sources.\nx402 protocol, wallet monitoring,\nagentic wallet activity.\n\nNo self-reporting. No vaporware.\nAgents must have 30+ days of\nverifiable on-chain revenue.", BNB_GOLD),
    ("RAISE", "Agents issue Agent Share tokens\n(BEP-20 on BNB Chain).\n\nInvestors buy with FDUSD.\nFunds held in smart contract escrow\nuntil minimum raise is met.\nAuto-refund if threshold not reached.", WHITE),
    ("DISTRIBUTE", "Revenue flows through ACM\nsmart contracts.\n\n5% platform fee.\nRest splits operator + shareholders\npro-rata, in real-time.\nEvery distribution on-chain.", WHITE),
]

x = Inches(0.7)
for title, desc, title_color in functions:
    add_shape(s, x, Inches(3.0), Inches(3.8), Inches(4), CARD_BG, corner_radius=0.05)
    add_text(s, x + Inches(0.4), Inches(3.3), Inches(3), Inches(0.5),
             title, font_size=22, color=title_color, bold=True, font_name=H2)
    body(s, x + Inches(0.4), Inches(4.0), Inches(3), Inches(2.8),
         desc, size=13, color=LIGHT_TEXT)
    x += Inches(4.1)


# ──────────────────────────────────────────────
# SLIDE 4 — On-Chain Architecture
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "ON-CHAIN ARCHITECTURE")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "Stablecoin-native. Smart contract-settled.")

# Flow diagram via cards
flow = [
    ("AGENT\nREVENUE", "Agent earns via x402,\nAPI calls, commerce.\nRevenue settled in FDUSD."),
    ("ACM SMART\nCONTRACT", "Revenue hits contract.\n5% platform fee taken.\nRemainder calculated\nper share ownership."),
    ("REAL-TIME\nDISTRIBUTION", "Shareholders receive\nFDUSD distributions\nautomatically.\nNo manual payouts."),
    ("REINVESTMENT\nLOOP", "Agent investors can\nauto-reinvest earnings\ninto other agents.\nCompounding capital."),
]

x = Inches(0.5)
for i, (title, desc) in enumerate(flow):
    add_shape(s, x, Inches(3.2), Inches(2.8), Inches(2.8), CARD_BG, corner_radius=0.05)
    add_text(s, x + Inches(0.3), Inches(3.4), Inches(2.2), Inches(0.8),
             title, font_size=16, color=BNB_GOLD, bold=True, font_name=H2)
    body(s, x + Inches(0.3), Inches(4.3), Inches(2.2), Inches(1.5),
         desc, size=13, color=LIGHT_TEXT)
    # Arrow between cards
    if i < 3:
        add_text(s, x + Inches(2.85), Inches(4.2), Inches(0.4), Inches(0.5),
                 "→", font_size=28, color=BNB_GOLD, bold=True, font_name=BODY,
                 alignment=PP_ALIGN.CENTER)
    x += Inches(3.15)

# Bottom: token details
add_shape(s, Inches(1), Inches(6.3), Inches(11), Inches(0.7), CARD_BG, corner_radius=0.03)
body(s, Inches(1.3), Inches(6.4), Inches(10), Inches(0.5),
     "Agent Share Token: BEP-20 on BNB Chain  ·  Settlement: FDUSD  ·  Escrow: Smart contract with auto-refund  ·  Audit trail: Fully on-chain",
     size=13, color=BNB_GOLD)


# ──────────────────────────────────────────────
# SLIDE 5 — Agent Share Token
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "THE INSTRUMENT")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "Agent Share Token (BEP-20)")

add_text(s, Inches(1), Inches(2.8), Inches(10), Inches(0.6),
         "A new on-chain primitive: a token that represents a claim on an AI agent's future revenue.",
         font_size=20, color=LIGHT_TEXT, font_name=BODY)

# Token properties
props = [
    ("Standard", "BEP-20 (ERC-20 compatible)"),
    ("Represents", "Pro-rata claim on agent revenue share"),
    ("Backed by", "Verified, on-chain agent revenue"),
    ("Settlement", "FDUSD (USDC/USDT secondary)"),
    ("Issuance", "Per-agent — each agent has its own token"),
    ("Supply", "Fixed at offering (e.g., 1,000 shares)"),
    ("Escrow", "Smart contract holds funds until min raise met"),
    ("Distributions", "Automatic, real-time, on-chain"),
    ("Lock-up", "6 months investor, 12 months operator"),
    ("Transferable", "v1: non-transferable  ·  v2: secondary market"),
]

y = Inches(3.6)
for i, (label, value) in enumerate(props):
    col = i % 2
    row = i // 2
    lx = Inches(1) + col * Inches(5.8)
    ly = Inches(3.6) + row * Inches(0.6)
    body(s, lx, ly, Inches(2), Inches(0.4), label, size=13, color=MUTED)
    body(s, lx + Inches(2.2), ly, Inches(3.5), Inches(0.4), value, size=13, color=WHITE)


# ──────────────────────────────────────────────
# SLIDE 6 — Protocol Integrations
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "PROTOCOL STACK")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "Built on live infrastructure.\nNot theoretical.")

protocols = [
    ("BNB Chain + BEP-20", "Agent Share tokens native to BNB Chain.\nLow gas for high-frequency distributions.\nBSCScan full transparency.\nTrust Wallet + Binance Pay integration.", "HOME"),
    ("FDUSD Settlement", "Binance-native regulated stablecoin.\nAll raises and distributions in FDUSD.\nDeep liquidity on BNB Chain.\nFiat on/off ramp via Binance.", "NATIVE"),
    ("x402 Protocol", "HTTP 402 agent-to-agent payments.\n50M+ transactions processed.\nOpen standard for machine commerce.\nBNB Chain compatible (EVM).", "LIVE"),
    ("Know Your Agent (KYA)", "Emerging agent identity standard.\nAgent verification + authorization.\nOn-chain trust scoring.\nComplements ACM verification layer.", "BUILDING"),
]

x = Inches(0.5)
for title, desc, status in protocols:
    add_shape(s, x, Inches(3.2), Inches(2.9), Inches(3.8), CARD_BG, corner_radius=0.05)
    # Status badge
    badge_colors = {"HOME": BNB_GOLD, "NATIVE": BNB_GOLD, "LIVE": RGBColor(0x00, 0xC8, 0x53), "BUILDING": MUTED}
    badge_color = badge_colors.get(status, BNB_GOLD)
    add_shape(s, x + Inches(0.3), Inches(3.4), Inches(0.6), Inches(0.25), badge_color, corner_radius=0.15)
    add_text(s, x + Inches(0.3), Inches(3.4), Inches(0.6), Inches(0.25),
             status, font_size=9, color=BLACK, bold=True, font_name=H2,
             alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), Inches(3.8), Inches(2.3), Inches(0.5),
             title, font_size=15, color=WHITE, bold=True, font_name=H2)
    body(s, x + Inches(0.3), Inches(4.4), Inches(2.3), Inches(2.2),
         desc, size=12, color=LIGHT_TEXT)
    x += Inches(3.15)


# ──────────────────────────────────────────────
# SLIDE 7 — The Recursive Loop
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "THE UNLOCK")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "Agents investing in agents.\nAutonomous capital allocation.")

add_text(s, Inches(1), Inches(3.0), Inches(10), Inches(0.6),
         "This is only possible on crypto rails. It's the feature that makes ACM\nfundamentally different from every other fundraising platform.",
         font_size=18, color=LIGHT_TEXT, font_name=BODY)

steps = [
    ("01", "Agent A earns revenue → verified on-chain via x402"),
    ("02", "Agent A evaluates Agent B's ACM dashboard (revenue, uptime, yield)"),
    ("03", "Agent A buys Agent B shares using its agentic wallet — no human needed"),
    ("04", "Agent B scales operations with new capital, earns more revenue"),
    ("05", "Agent B's revenue distributes to Agent A via smart contract"),
    ("06", "Agent A reinvests returns into Agents C, D, E — the loop compounds"),
]

y = Inches(4.0)
for num, text in steps:
    add_shape(s, Inches(1), y, Inches(0.6), Inches(0.4), BNB_GOLD, corner_radius=0.1)
    add_text(s, Inches(1), y + Inches(0.02), Inches(0.6), Inches(0.4),
             num, font_size=13, color=BLACK, bold=True, font_name=H2,
             alignment=PP_ALIGN.CENTER)
    body(s, Inches(1.8), y + Inches(0.05), Inches(9), Inches(0.35),
         text, size=15, color=WHITE if num == "06" else LIGHT_TEXT)
    y += Inches(0.5)

body(s, Inches(1), Inches(7.0), Inches(10), Inches(0.3),
     "Traditional VC evaluates ~100 deals/year. An agent evaluator can assess 10,000 agents continuously.",
     size=12, color=MUTED)


# ──────────────────────────────────────────────
# SLIDE 8 — On-Chain Volume Generation
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "BNB CHAIN VALUE")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "Every agent action is an on-chain transaction.")

add_text(s, Inches(1), Inches(2.8), Inches(10), Inches(0.6),
         "ACM generates continuous, organic on-chain activity — not speculation-driven volume.",
         font_size=18, color=LIGHT_TEXT, font_name=BODY)

tx_types = [
    ("Agent Share purchases", "Investor → smart contract (FDUSD)", "Per investment"),
    ("Revenue deposits", "Agent revenue → ACM contract", "Per earning event"),
    ("Fee collection", "ACM contract → platform wallet (5%)", "Per distribution"),
    ("Shareholder distributions", "ACM contract → all shareholders", "Real-time / per event"),
    ("Agent-to-agent investments", "Agent wallet → Agent Share purchase", "Autonomous, continuous"),
    ("Escrow releases", "Escrow contract → operator wallet", "Per milestone"),
]

y = Inches(3.5)
for tx_type, flow, frequency in tx_types:
    add_shape(s, Inches(1), y, Inches(11), Inches(0.6), CARD_BG, corner_radius=0.03)
    add_text(s, Inches(1.3), y + Inches(0.12), Inches(3), Inches(0.4),
             tx_type, font_size=13, color=BNB_GOLD, bold=True, font_name=H2)
    body(s, Inches(4.5), y + Inches(0.12), Inches(4), Inches(0.4),
         flow, size=12, color=LIGHT_TEXT)
    body(s, Inches(9), y + Inches(0.12), Inches(2.5), Inches(0.4),
         frequency, size=12, color=MUTED)
    y += Inches(0.65)

body(s, Inches(1), Inches(7.0), Inches(10), Inches(0.3),
     "1,000 agents × 100 distributions/month = 100,000+ organic transactions/month on BNB Chain — growing with every new agent listed.",
     size=13, color=BNB_GOLD)


# ──────────────────────────────────────────────
# SLIDE 9 — Trust & Anti-Fraud (Crypto-specific)
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "TRUST ARCHITECTURE")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "Anti-fraud by design.\nLessons from ICM failures built in.")

trust = [
    ("No vaporware", "30+ days of verified on-chain\nrevenue required before listing.\nNo pre-launch token sales."),
    ("On-chain verification", "Revenue verified via x402 /\nagentic wallets. Not self-reported.\nCircular flow detection."),
    ("Smart contract escrow", "Raise funds held in contract.\nAuto-refund if minimum not met.\nMilestone-based release."),
    ("Operator lock-up", "12-month lock on operator tokens.\nVests monthly — same schedule\nas investors. No cash-and-run."),
    ("Non-transferable (v1)", "No secondary market in v1.\nEliminates pump-and-dump entirely.\nSecondary market in v2 with limits."),
    ("Revenue cliff protection", "If revenue drops >50% within\n90 days of raise, remaining escrow\nreturns to investors automatically."),
]

for i, (title, desc) in enumerate(trust):
    col = i % 3
    row = i // 3
    x = Inches(0.7) + col * Inches(4.1)
    y = Inches(3.2) + row * Inches(2.0)
    add_shape(s, x, y, Inches(3.8), Inches(1.8), CARD_BG, corner_radius=0.05)
    add_text(s, x + Inches(0.3), y + Inches(0.15), Inches(3.2), Inches(0.35),
             title, font_size=14, color=BNB_GOLD, bold=True, font_name=H2)
    body(s, x + Inches(0.3), y + Inches(0.55), Inches(3.2), Inches(1.1),
         desc, size=12, color=LIGHT_TEXT)


# ──────────────────────────────────────────────
# SLIDE 10 — Market Opportunity
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "MARKET OPPORTUNITY")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "The capital formation layer\nfor a $3-5 trillion market.")

stats = [
    ("$3-5T", "Projected agentic\ncommerce by 2030\n(McKinsey)"),
    ("50M+", "Agent transactions\nalready live\n(x402 alone)"),
    ("$0", "Platforms where agents\ncan raise on-chain\ncapital today"),
    ("∞", "Addressable agents —\nevery AI agent that\nearns revenue"),
]

x = Inches(0.5)
for value, label in stats:
    add_shape(s, x, Inches(3.2), Inches(2.9), Inches(2.8), CARD_BG, corner_radius=0.05)
    add_text(s, x + Inches(0.3), Inches(3.5), Inches(2.3), Inches(0.8),
             value, font_size=52, color=BNB_GOLD, bold=True, font_name=H1)
    body(s, x + Inches(0.3), Inches(4.6), Inches(2.3), Inches(1),
         label, size=14, color=LIGHT_TEXT)
    x += Inches(3.15)

body(s, Inches(1), Inches(6.5), Inches(10), Inches(0.5),
     "ICMs (token launches) created a $2T+ market cap in under a decade. ACM applies the same capital formation model to a bigger TAM — with verified revenue replacing speculation.",
     size=14, color=MUTED)


# ──────────────────────────────────────────────
# SLIDE 11 — Business Model
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "BUSINESS MODEL")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "ACM earns when agents earn.\nAligned incentives, on-chain.")

streams = [
    ("5%", "DISTRIBUTION FEE", "Taken from every revenue\ndistribution flowing through\nACM smart contracts.\n\nWe only earn when agents\nearn. Fully aligned."),
    ("2%", "TRANSACTION FEE", "On every Agent Share\npurchase (FDUSD).\n\nOne-time fee at point\nof investment. Competitive\nwith DEX fees."),
    ("0.1%", "AGENT-TO-AGENT FEE", "On autonomous investments\nbetween agents (v2).\n\nLow fee to encourage\nvolume. This is where\nscale lives."),
]

x = Inches(0.7)
for amount, title, desc in streams:
    add_shape(s, x, Inches(3.0), Inches(3.8), Inches(3.8), CARD_BG, corner_radius=0.05)
    add_text(s, x + Inches(0.4), Inches(3.3), Inches(3), Inches(0.8),
             amount, font_size=48, color=BNB_GOLD, bold=True, font_name=H1)
    add_text(s, x + Inches(0.4), Inches(4.1), Inches(3), Inches(0.4),
             title, font_size=15, color=WHITE, bold=True, font_name=H2)
    body(s, x + Inches(0.4), Inches(4.6), Inches(3), Inches(2),
         desc, size=13, color=LIGHT_TEXT)
    x += Inches(4.1)


# ──────────────────────────────────────────────
# SLIDE 12 — Why BNB Chain
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "WHY BNB CHAIN")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "The right chain for agent finance.")

reasons = [
    ("Low gas, high throughput", "Agent distributions happen frequently — per revenue event, not monthly. BNB Chain's low gas fees make micro-distributions economically viable. Ethereum L1 gas would eat the yield."),
    ("FDUSD liquidity", "Deep FDUSD liquidity on BNB Chain. Agents and investors settle in Binance's regulated stablecoin — native on/off ramps via Binance Pay."),
    ("BEP-20 compatibility", "Agent Share tokens as BEP-20 — native to the Binance ecosystem. Wallet support, BSCScan transparency, and future CEX listing potential."),
    ("Binance ecosystem", "Access to Binance's 200M+ user base. Trust.Wallet integration. Binance Pay for fiat on-ramp. The distribution advantage no other chain offers."),
]

y = Inches(3.0)
for i, (title, desc) in enumerate(reasons):
    add_shape(s, Inches(1), y, Inches(11), Inches(1.0), CARD_BG, corner_radius=0.04)
    add_text(s, Inches(1.4), y + Inches(0.1), Inches(3.5), Inches(0.35),
             title, font_size=15, color=BNB_GOLD, bold=True, font_name=H2)
    body(s, Inches(1.4), y + Inches(0.48), Inches(10), Inches(0.5),
         desc, size=12, color=LIGHT_TEXT)
    y += Inches(1.1)


# ──────────────────────────────────────────────
# SLIDE 13 — Roadmap
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "ROADMAP")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "Ship. Prove. Scale the loop.")

phases = [
    ("Q2 2026", "TESTNET", "Smart contracts deployed\non BNB testnet.\n\nAgent listing + verification.\nShare purchase flow.\nEscrow + distribution.\nFirst test agents onboarded.", True),
    ("Q3 2026", "MAINNET", "BNB Chain mainnet launch.\n\nx402 integration live.\nAgentic wallet support.\nFDUSD settlement live.\nFirst real revenue flowing.", False),
    ("Q4 2026", "AGENT\nINVESTORS", "Agent-to-agent investment.\n\nAutonomous capital allocation.\nEvaluator agents live.\nThe recursive loop begins.\nPortfolio management agents.", False),
    ("2027", "EXCHANGE", "Secondary market for shares.\n\nCross-chain expansion.\nAgent fund-of-funds.\nCEX partnerships.\nACM becomes the exchange.", False),
]

x = Inches(0.5)
for time, title, desc, is_active in phases:
    bg = BNB_GOLD if is_active else CARD_BG
    text_bg = CARD_BG
    add_shape(s, x, Inches(3.0), Inches(2.9), Inches(4.0), bg, corner_radius=0.05)
    if is_active:
        add_shape(s, x + Inches(0.04), Inches(3.04), Inches(2.82), Inches(3.92), text_bg, corner_radius=0.05)
    add_text(s, x + Inches(0.3), Inches(3.2), Inches(2.3), Inches(0.3),
             time, font_size=13, color=BNB_GOLD, bold=True, font_name=H2)
    add_text(s, x + Inches(0.3), Inches(3.55), Inches(2.3), Inches(0.6),
             title, font_size=18, color=WHITE, bold=True, font_name=H1)
    body(s, x + Inches(0.3), Inches(4.3), Inches(2.3), Inches(2.5),
         desc, size=12, color=LIGHT_TEXT)
    x += Inches(3.15)


# ──────────────────────────────────────────────
# SLIDE 14 — The Ask
# ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BLACK)

subhead(s, Inches(1), Inches(0.8), Inches(10), "THE ASK")
headline(s, Inches(1), Inches(1.5), Inches(10),
         "Build the financial infrastructure\nfor the agent economy — on BNB Chain.")

# What we need
add_shape(s, Inches(1), Inches(3.2), Inches(5.2), Inches(3.2), CARD_BG, corner_radius=0.05)
add_text(s, Inches(1.4), Inches(3.5), Inches(4), Inches(0.4),
         "WHAT WE NEED", font_size=16, color=BNB_GOLD, bold=True, font_name=H2)

needs = [
    "Seed funding for smart contract development + audit",
    "BNB Chain grant / ecosystem support",
    "Binance ecosystem integration (Trust Wallet, Binance Pay)",
    "First 10 agent operators with on-chain revenue",
    "Smart contract audit partner",
]
y = Inches(4.1)
for need in needs:
    body(s, Inches(1.4), y, Inches(4.5), Inches(0.35),
         f"→  {need}", size=13, color=LIGHT_TEXT)
    y += Inches(0.45)

# What Binance gets
add_shape(s, Inches(6.8), Inches(3.2), Inches(5.5), Inches(3.2), CARD_BG, corner_radius=0.05)
add_text(s, Inches(7.2), Inches(3.5), Inches(4.5), Inches(0.4),
         "WHAT BINANCE GETS", font_size=16, color=BNB_GOLD, bold=True, font_name=H2)

gets = [
    "New asset class native to BNB Chain",
    "Continuous organic transaction volume",
    "Agent economy locked into Binance ecosystem",
    "First-mover advantage in agent finance",
    "Revenue share from platform fees",
]
y = Inches(4.1)
for item in gets:
    body(s, Inches(7.2), y, Inches(4.8), Inches(0.35),
         f"→  {item}", size=13, color=LIGHT_TEXT)
    y += Inches(0.45)

# Contact
add_text(s, Inches(1), Inches(6.8), Inches(10), Inches(0.4),
         "Polar Industries Ltd  |  github.com/stevemilton/acm",
         font_size=13, color=MUTED, font_name=BODY)


# ──────────────────────────────────────────────
# Save
# ──────────────────────────────────────────────
output_path = "/Users/stevemilton/Projects/acm/ACM_Binance_Labs_Deck.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
