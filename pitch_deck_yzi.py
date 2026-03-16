#!/usr/bin/env python3
"""Generate ACM pitch deck for YZi Labs — tight, 8 slides, BNB-native."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BLACK = RGBColor(0x08, 0x08, 0x08)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF0, 0xB9, 0x0B)
CARD = RGBColor(0x18, 0x18, 0x18)
MUTED = RGBColor(0x77, 0x77, 0x77)
LIGHT = RGBColor(0xCC, 0xCC, 0xCC)

H1 = "Montserrat"
H2 = "Play"
BODY = "Inter"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = BLACK


def box(slide, l, t, w, h, color=CARD, r=0.05):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.adjustments[0] = r
    return s


def txt(slide, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT, f=BODY, font_name=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = font_name if font_name else f
    p.alignment = a
    return tb


# ── SLIDE 1 — Title ──
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

txt(s, Inches(1), Inches(1), Inches(10), Inches(0.5),
    "FOR YZI LABS", 16, GOLD, True, font_name=H2)
txt(s, Inches(1), Inches(2), Inches(11), Inches(1.5),
    "AGENT\nCAPITAL MARKETS", 64, WHITE, True, font_name=H1)
txt(s, Inches(1), Inches(4.2), Inches(8), Inches(0.8),
    "The exchange where AI agents raise capital\nand distribute revenue — on BNB Chain.",
    26, LIGHT, font_name=BODY)
txt(s, Inches(1), Inches(6.2), Inches(6), Inches(0.4),
    "FDUSD settlement  ·  BEP-20 tokens  ·  Agent-to-agent investment",
    14, GOLD, True, font_name=H2)


# ── SLIDE 2 — Problem / Opportunity ──
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

txt(s, Inches(1), Inches(0.8), Inches(10), Inches(0.4),
    "THE OPPORTUNITY", 14, GOLD, True, font_name=H2)
txt(s, Inches(1), Inches(1.5), Inches(10), Inches(1.2),
    "Agents earn revenue.\nNo one lets them raise capital.", 40, WHITE, True, font_name=H1)

stats = [
    ("$3-5T", "agentic commerce\nby 2030 (McKinsey)"),
    ("50M+", "agent transactions\nalready live (x402)"),
    ("ZERO", "platforms for agents\nto raise capital"),
]
x = Inches(1)
for val, label in stats:
    box(s, x, Inches(3.5), Inches(3.5), Inches(2.5))
    txt(s, x + Inches(0.4), Inches(3.8), Inches(2.7), Inches(0.8),
        val, 48, GOLD, True, font_name=H1)
    txt(s, x + Inches(0.4), Inches(4.8), Inches(2.7), Inches(1),
        label, 16, LIGHT, font_name=BODY)
    x += Inches(3.8)

txt(s, Inches(1), Inches(6.5), Inches(10), Inches(0.4),
    "Virtuals Protocol proved demand ($5B mcap) — but it's speculative. ACM is the verified version.",
    14, MUTED, font_name=BODY)


# ── SLIDE 3 — What ACM Does ──
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

txt(s, Inches(1), Inches(0.8), Inches(10), Inches(0.4),
    "HOW IT WORKS", 14, GOLD, True, font_name=H2)
txt(s, Inches(1), Inches(1.5), Inches(10), Inches(1),
    "Verify. Raise. Distribute.", 44, WHITE, True, font_name=H1)

funcs = [
    ("VERIFY", "Agents connect revenue sources\n(x402, Stripe, on-chain).\n\n30+ days of verified revenue\nrequired. No vaporware.\nReal-time performance dashboard."),
    ("RAISE", "Agents issue BEP-20 revenue\nshare tokens on BNB Chain.\n\nInvestors buy with FDUSD.\nSmart contract escrow.\nAuto-refund if min not met."),
    ("DISTRIBUTE", "Revenue flows through ACM\nsmart contracts automatically.\n\n5% platform fee. Rest splits\noperator + shareholders pro-rata.\nReal-time, on-chain."),
]
x = Inches(0.7)
for title, desc in funcs:
    box(s, x, Inches(3.0), Inches(3.8), Inches(4))
    txt(s, x + Inches(0.4), Inches(3.3), Inches(3), Inches(0.5),
        title, 22, GOLD, True, font_name=H2)
    txt(s, x + Inches(0.4), Inches(4.0), Inches(3), Inches(2.8),
        desc, 14, LIGHT, font_name=BODY)
    x += Inches(4.1)


# ── SLIDE 4 — The Unlock ──
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

txt(s, Inches(1), Inches(0.8), Inches(10), Inches(0.4),
    "THE UNLOCK", 14, GOLD, True, font_name=H2)
txt(s, Inches(1), Inches(1.5), Inches(10), Inches(1.2),
    "Agents investing in agents.", 44, WHITE, True, font_name=H1)

txt(s, Inches(1), Inches(3.0), Inches(9), Inches(0.5),
    "Traditional VC evaluates ~100 deals/year. An evaluator agent can assess 10,000 continuously.",
    18, LIGHT, font_name=BODY)

loop = [
    "Agent A earns revenue → verified on-chain",
    "Agent A evaluates Agent B's performance on ACM",
    "Agent A buys Agent B shares via agentic wallet — no human needed",
    "Agent B earns more → revenue distributes back to Agent A",
    "Agent A reinvests into C, D, E → the loop compounds",
]
y = Inches(4.0)
for i, step in enumerate(loop):
    box(s, Inches(1), y, Inches(0.55), Inches(0.4), GOLD, r=0.15)
    txt(s, Inches(1), y + Inches(0.02), Inches(0.55), Inches(0.4),
        f"0{i+1}", 12, BLACK, True, a=PP_ALIGN.CENTER, font_name=H2)
    txt(s, Inches(1.8), y + Inches(0.04), Inches(9), Inches(0.35),
        step, 16, GOLD if i == 4 else LIGHT, i == 4, font_name=BODY)
    y += Inches(0.52)

txt(s, Inches(1), Inches(6.8), Inches(10), Inches(0.3),
    "Only possible on crypto rails. Fiat requires human approval for every transaction.",
    12, MUTED, font_name=BODY)


# ── SLIDE 5 — Why BNB Chain ──
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

txt(s, Inches(1), Inches(0.8), Inches(10), Inches(0.4),
    "WHY BNB CHAIN", 14, GOLD, True, font_name=H2)
txt(s, Inches(1), Inches(1.5), Inches(10), Inches(1),
    "Built for Binance. Not ported.", 44, WHITE, True, font_name=H1)

reasons = [
    ("FDUSD native", "All settlement in Binance's regulated stablecoin.\nDrives FDUSD demand with every raise and distribution."),
    ("Low gas = viable micro-distributions", "Agent distributions happen per revenue event, not monthly.\nBNB Chain gas makes this economically viable."),
    ("200M+ user distribution", "Trust Wallet, Binance Pay, BSCScan.\nNo other chain offers this investor on-ramp."),
    ("Agent-native primitives", "ERC-8004 agent identities + BAP-578 Non-Fungible Agents\nalready live on BNB Chain. ACM builds on top."),
]
y = Inches(3.0)
for title, desc in reasons:
    box(s, Inches(1), y, Inches(11), Inches(0.9))
    txt(s, Inches(1.4), y + Inches(0.1), Inches(4), Inches(0.3),
        title, 15, GOLD, True, font_name=H2)
    txt(s, Inches(1.4), y + Inches(0.45), Inches(10), Inches(0.4),
        desc, 12, LIGHT, font_name=BODY)
    y += Inches(1.0)


# ── SLIDE 6 — Business Model ──
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

txt(s, Inches(1), Inches(0.8), Inches(10), Inches(0.4),
    "BUSINESS MODEL", 14, GOLD, True, font_name=H2)
txt(s, Inches(1), Inches(1.5), Inches(10), Inches(1),
    "ACM earns when agents earn.", 44, WHITE, True, font_name=H1)

fees = [
    ("5%", "distribution fee", "On all revenue flowing\nthrough smart contracts"),
    ("2%", "transaction fee", "On every Agent Share\npurchase (FDUSD)"),
    ("0.1%", "agent-to-agent", "On autonomous investments\nbetween agents (v2)"),
]
x = Inches(0.7)
for amt, label, desc in fees:
    box(s, x, Inches(3.0), Inches(3.8), Inches(2.2))
    txt(s, x + Inches(0.4), Inches(3.2), Inches(3), Inches(0.8),
        amt, 52, GOLD, True, font_name=H1)
    txt(s, x + Inches(0.4), Inches(4.0), Inches(3), Inches(0.3),
        label, 16, WHITE, True, font_name=H2)
    txt(s, x + Inches(0.4), Inches(4.4), Inches(3), Inches(0.7),
        desc, 13, LIGHT, font_name=BODY)
    x += Inches(4.1)

# Scale table
box(s, Inches(1), Inches(5.6), Inches(11), Inches(1.5))
txt(s, Inches(1.4), Inches(5.7), Inches(3), Inches(0.3),
    "SCALE PROJECTIONS", 12, MUTED, True, font_name=H2)

headers = ["", "100 agents", "500 agents", "1,000 agents"]
row1 = ["Monthly distributions", "$200K", "$1M", "$2M"]
row2 = ["Annual platform revenue", "$170K", "$850K", "$1.7M"]

hx = Inches(1.4)
for h in headers:
    txt(s, hx, Inches(6.05), Inches(2.3), Inches(0.3),
        h, 11, MUTED, True, font_name=H2)
    hx += Inches(2.5)

hx = Inches(1.4)
for v in row1:
    txt(s, hx, Inches(6.35), Inches(2.3), Inches(0.3),
        v, 12, LIGHT, font_name=BODY)
    hx += Inches(2.5)

hx = Inches(1.4)
for v in row2:
    txt(s, hx, Inches(6.65), Inches(2.3), Inches(0.3),
        v, 12, GOLD if v.startswith("$") else LIGHT, v.startswith("$"), font_name=BODY)
    hx += Inches(2.5)


# ── SLIDE 7 — Competitors ──
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

txt(s, Inches(1), Inches(0.8), Inches(10), Inches(0.4),
    "COMPETITIVE POSITION", 14, GOLD, True, font_name=H2)
txt(s, Inches(1), Inches(1.5), Inches(10), Inches(1),
    "Everyone else built the casino.\nWe're building the exchange.", 40, WHITE, True, font_name=H1)

comps = [
    ("Virtuals Protocol", "Base", "$5B peak mcap", "Speculative IAOs, bonding curves,\nno verified revenue, no protections"),
    ("ElizaOS / ai16z", "Solana", "$2.5B peak mcap", "AI-managed fund — one token,\nnot individual agent investment"),
    ("AgentStocks", "—", "Early", "Capital TO agents for trading,\nnot investment IN agents"),
    ("ACM", "BNB Chain", "Pre-launch", "Verified revenue, structured shares,\nescrow, lock-ups, agent-to-agent"),
]

y = Inches(3.2)
for name, chain, scale, desc in comps:
    is_acm = name == "ACM"
    box(s, Inches(1), y, Inches(11), Inches(0.85), GOLD if is_acm else CARD)
    if is_acm:
        box(s, Inches(1.03), y + Inches(0.03), Inches(10.94), Inches(0.79), CARD)
    txt(s, Inches(1.4), y + Inches(0.12), Inches(2.5), Inches(0.3),
        name, 15, GOLD if is_acm else WHITE, True, font_name=H2)
    txt(s, Inches(4), y + Inches(0.12), Inches(1.5), Inches(0.3),
        chain, 12, GOLD if is_acm else MUTED, font_name=BODY)
    txt(s, Inches(5.5), y + Inches(0.12), Inches(1.8), Inches(0.3),
        scale, 12, GOLD if is_acm else MUTED, font_name=BODY)
    txt(s, Inches(7.5), y + Inches(0.12), Inches(4), Inches(0.65),
        desc, 12, GOLD if is_acm else LIGHT, font_name=BODY)
    y += Inches(0.95)


# ── SLIDE 8 — The Ask ──
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

txt(s, Inches(1), Inches(0.8), Inches(10), Inches(0.4),
    "THE ASK", 14, GOLD, True, font_name=H2)
txt(s, Inches(1), Inches(1.5), Inches(10), Inches(1),
    "The agent economy needs an exchange.\nLet's build it on BNB Chain.", 40, WHITE, True, font_name=H1)

# Two columns
box(s, Inches(1), Inches(3.2), Inches(5.2), Inches(3))
txt(s, Inches(1.4), Inches(3.4), Inches(4), Inches(0.4),
    "WHAT WE NEED", 16, GOLD, True, font_name=H2)
needs = [
    "Seed funding — smart contracts + audit",
    "BNB Chain ecosystem support",
    "Trust Wallet + Binance Pay integration",
    "10 launch agents with on-chain revenue",
]
y = Inches(4.0)
for n in needs:
    txt(s, Inches(1.4), y, Inches(4.5), Inches(0.3),
        f"→  {n}", 14, LIGHT, font_name=BODY)
    y += Inches(0.45)

box(s, Inches(6.8), Inches(3.2), Inches(5.5), Inches(3))
txt(s, Inches(7.2), Inches(3.4), Inches(4.5), Inches(0.4),
    "WHAT YZI LABS GETS", 16, GOLD, True, font_name=H2)
gets = [
    "New asset class native to BNB Chain",
    "Continuous organic on-chain volume",
    "FDUSD demand from every transaction",
    "First-mover in agent capital markets",
]
y = Inches(4.0)
for g in gets:
    txt(s, Inches(7.2), y, Inches(4.8), Inches(0.3),
        f"→  {g}", 14, LIGHT, font_name=BODY)
    y += Inches(0.45)

# Timeline bar
box(s, Inches(1), Inches(6.5), Inches(11), Inches(0.6))
phases = [
    ("Q2 2026", "Testnet"),
    ("Q3 2026", "BNB Mainnet"),
    ("Q4 2026", "Agent-to-Agent"),
    ("2027", "Exchange"),
]
x = Inches(1.3)
for time, phase in phases:
    txt(s, x, Inches(6.55), Inches(1.2), Inches(0.25),
        time, 11, GOLD, True, font_name=H2)
    txt(s, x + Inches(1.3), Inches(6.55), Inches(1.5), Inches(0.25),
        phase, 11, LIGHT, font_name=BODY)
    x += Inches(2.7)

txt(s, Inches(1), Inches(7.2), Inches(10), Inches(0.3),
    "Polar Industries Ltd  ·  github.com/stevemilton/acm",
    12, MUTED, font_name=BODY)


# ── Save ──
out = "/Users/stevemilton/Projects/acm/ACM_YZi_Labs_Deck.pptx"
prs.save(out)
print(f"Saved to {out}")
